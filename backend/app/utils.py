import os
import csv
import json
import logging
import subprocess
import platform
import psutil
import ipaddress
import socket
import threading
from typing import Optional
from urllib.parse import urljoin, urlparse
from pypdf import PdfReader
from sqlalchemy.orm import Session
from app.config import settings

logger = logging.getLogger(__name__)
_whisper_model = None
_whisper_model_lock = threading.Lock()
_whisper_inference_lock = threading.Lock()


def _validate_public_url(url: str) -> None:
    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"} or not parsed.hostname or parsed.username or parsed.password:
        raise ValueError("Only public http/https URLs are supported.")
    if parsed.port not in {None, 80, 443}:
        raise ValueError("Only standard web ports 80 and 443 are supported.")
    try:
        addresses = {item[4][0] for item in socket.getaddrinfo(parsed.hostname, parsed.port or (443 if parsed.scheme == "https" else 80), type=socket.SOCK_STREAM)}
    except socket.gaierror as exc:
        raise ValueError("The website hostname could not be resolved.") from exc
    for address in addresses:
        ip = ipaddress.ip_address(address)
        if ip.is_private or ip.is_loopback or ip.is_link_local or ip.is_reserved or ip.is_multicast or ip.is_unspecified:
            raise ValueError("Private, local, reserved, and metadata-network URLs are blocked.")


def _safe_public_get(url: str, headers: dict, timeout: int = 15):
    import requests as _r
    current = url
    for _ in range(6):
        _validate_public_url(current)
        response = _r.get(current, headers=headers, timeout=timeout, allow_redirects=False)
        if response.is_redirect or response.is_permanent_redirect:
            location = response.headers.get("location")
            if not location:
                break
            current = urljoin(current, location)
            continue
        return response
    raise ValueError("The website redirected too many times.")


# Speech recognition quality is very uneven across Whisper sizes: `base` is fast
# and accurate enough for English, but it renders Indic speech as nonsense
# romanisation ("बैटरी कितनी है" -> "Bachelor is so much."). Those languages need
# `small`, which costs a few extra seconds but actually returns their script.
_WHISPER_ENGLISH_MODEL = os.getenv("UPLOAD_WHISPER_MODEL", "base")
_WHISPER_MULTILINGUAL_MODEL = os.getenv("UPLOAD_WHISPER_MODEL_MULTILINGUAL", "small")
# Only for the provisional readings shown while somebody is still speaking.
# Measured at 1.71s against base's 2.49s on the same two-second clip.
_WHISPER_LIVE_MODEL = os.getenv("LIVE_WHISPER_MODEL", "tiny.en")
_whisper_models: dict = {}


def _whisper_model_name_for(language: str) -> str:
    normalized = (language or "auto").lower().split("-")[0]
    if normalized in {"en", ""}:
        return _WHISPER_ENGLISH_MODEL
    # "auto" may well be a non-English utterance, so use the accurate model.
    return _WHISPER_MULTILINGUAL_MODEL


def _get_whisper_model(model_name: str):
    """Load and cache a Whisper model by size."""
    global _whisper_model

    cached = _whisper_models.get(model_name)
    if cached is not None:
        return cached

    from faster_whisper import WhisperModel

    with _whisper_model_lock:
        cached = _whisper_models.get(model_name)
        if cached is not None:
            return cached
        download_root = os.path.join(settings.DATA_DIR, "models", "faster-whisper")
        os.makedirs(download_root, exist_ok=True)
        # cpu_threads matters far more than it looks. Left unset, loading the
        # base model on a sixteen-thread machine took 15.9 seconds; told how
        # many threads it has, 1.4. Transcription itself is about two seconds
        # either way, so nearly all of the wait a person felt on their first
        # dictation was this - the model being loaded slowly, once.
        #
        # Read from the machine rather than fixed, and capped: past about
        # eight threads there is nothing more to gain here, and taking every
        # core would fight with whatever model is answering the chat.
        try:
            import os as _os
            threads = min(8, max(1, (_os.cpu_count() or 4)))
        except Exception:
            threads = 4

        # The card if it can really be driven, the processor otherwise.
        #
        # device was hardcoded to "cpu", which was right while nothing knew
        # how to tell a card that is present from one that can be used.
        # gpu_speech answers that by trying, and falls back on its own, so
        # this cannot end up asking for a card that will raise.
        try:
            from app import gpu_speech

            device, compute = gpu_speech.device_and_compute()
        except Exception:  # noqa: BLE001 - never stop speech from loading
            device, compute = "cpu", "int8"
        device = os.getenv("UPLOAD_WHISPER_DEVICE", device)

        model = WhisperModel(
            model_name,
            device=device,
            compute_type=compute if device == "cuda" else "int8",
            download_root=download_root,
            cpu_threads=int(os.getenv("UPLOAD_WHISPER_THREADS", threads)),
        )
        _whisper_models[model_name] = model
        # Keep the historical single-model global pointing at something valid.
        _whisper_model = model
        return model


def warm_up_speech_models() -> None:
    """Load the speech models ahead of first use.

    Called in the background at startup so the first spoken sentence is not
    delayed by a cold model load, which otherwise takes tens of seconds. The
    fast English model is loaded first, then the multilingual one.
    """
    for model_name in dict.fromkeys((_WHISPER_ENGLISH_MODEL, _WHISPER_MULTILINGUAL_MODEL)):
        try:
            _get_whisper_model(model_name)
            logger.info(f"Speech recognition model '{model_name}' is warm.")
        except Exception as exc:  # noqa: BLE001 - warm-up is best effort
            logger.warning(f"Speech model '{model_name}' warm-up skipped: {exc}")


#: Why the last transcription produced nothing, or None if it succeeded.
#:
#: Failures here used to vanish into a logger.warning and an empty string,
#: and the interface turned that silence into "No speech was recognized.
#: Check microphone permission" - blaming the microphone for something that
#: had nothing to do with it. Several callers depend on the empty string, so
#: the reason is recorded beside it rather than raised.
_last_transcription_error: Optional[str] = None


def last_transcription_error() -> Optional[str]:
    """The reason the last transcription returned nothing, if it failed."""
    return _last_transcription_error


def _transcribe_local_media(file_path: str, language: str = "auto",
                            live: bool = False) -> str:
    """Turn a recording into text.

    `live` is for dictation while somebody is still speaking, where the answer
    is provisional and will be asked for again a second later.

    MEASURED ON THIS MACHINE, on a two-second recording, fastest of three:

        tiny.en   1.71s
        base      2.49s
        small     2.52s

    So the model is the only lever that moves. Greedy decoding instead of a
    beam of two, and dropping the voice-activity filter, were expected to help
    and did not: 2.57s against 2.45s, which is noise. They are kept because
    they are the right shape for a throwaway reading, not because they are
    worth anything measurable here - and this comment says so rather than
    claiming a speedup that was not observed.

    Most of those seconds are fixed cost - loading the audio and one encoder
    pass - not decoding, which is why a two-second clip and a longer one cost
    roughly the same. That is also why the interface asks again on a timer and
    keeps only one request in flight, instead of trying to make each call
    quick.

    tiny.en is less accurate, and on the sample above it misheard the sentence
    outright. That is acceptable only because a full-quality pass replaces the
    text the moment recording stops. It would not be acceptable as the answer.

    condition_on_previous_text stays False in both. In streaming that is not
    an accuracy setting but a safety one: one bad chunk conditions every chunk
    after it and the transcript drifts somewhere else entirely.

    THE BEAM, MEASURED RATHER THAN GUESSED

    Reported as mishearing. On the same recording, base, fastest of three:

        beam 1   3.28s   "Do you want to do it?"
        beam 2   3.24s   "Do you want to sing?"
        beam 5   3.13s   "Do you want me to sing a song?"

    The widest search was both the most accurate and, within noise, the
    fastest - a wider beam finds a confident path sooner and stops
    re-deciding. The final pass was on 2 and is on 5; that is free.

    NO GPU HERE, AND WHY

    ctranslate2 reports one CUDA device on this machine, so the obvious
    move looked like device="cuda". Trying it fails:

        RuntimeError: Library cublas64_12.dll is not found or cannot be
        loaded

    Seeing the card is not the same as having the CUDA runtime to drive it,
    and cuBLAS and cuDNN are a few hundred megabytes this app does not ship.
    So this stays on the processor, and roughly three seconds is what a
    processor costs. Milliseconds need either those libraries installed or
    a hosted service, and both are choices somebody should make knowingly.
    """
    # Local-first and private by default. No hosted speech API is called here.
    global _last_transcription_error
    _last_transcription_error = None
    try:
        wanted = _whisper_model_name_for(language)
        if live:
            wanted = _WHISPER_LIVE_MODEL
        model = _get_whisper_model(wanted)
        selected_language = (language or "auto").lower().split("-")[0]
        whisper_language = None if selected_language in {"", "auto"} else selected_language
        with _whisper_inference_lock:
            segments, _ = model.transcribe(
                file_path,
                language=whisper_language,
                task="transcribe",
                beam_size=1 if live else 5,
                vad_filter=not live,
                condition_on_previous_text=False,
            )
        res = " ".join(segment.text.strip() for segment in segments if segment.text.strip())
        if res:
            return res
        # Ran fine and heard nothing. That is a real answer, not a fault, and
        # it is the one case where "check your microphone" is fair advice.
        _last_transcription_error = None
    except Exception as whisper_err:
        _last_transcription_error = "%s: %s" % (type(whisper_err).__name__, whisper_err)
        logger.warning("Local faster-whisper transcription failed: %s", whisper_err,
                       exc_info=True)
    return ""

# --- Ingestion File Parsers ---
def parse_file_content(file_path: str, file_type: str) -> str:
    """Parse text contents based on file extensions."""
    file_type = file_type.lower()

    if file_type == "pdf":
        try:
            reader = PdfReader(file_path)
            text_parts = []
            for page_number, page in enumerate(reader.pages, start=1):
                text = page.extract_text()
                if text:
                    text_parts.append(f"[Page {page_number}]\n{text}")
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error parsing PDF file {file_path}: {e}")
            raise ValueError(f"Could not parse PDF content: {str(e)}")
            
    elif file_type == "csv":
        try:
            formatted_rows = []
            with open(file_path, mode="r", encoding="utf-8-sig") as f:
                reader = csv.reader(f)
                for row_index, row in enumerate(reader, start=1):
                    cells = []
                    for column_index, value in enumerate(row, start=1):
                        column_name = ""
                        number = column_index
                        while number:
                            number, remainder = divmod(number - 1, 26)
                            column_name = chr(65 + remainder) + column_name
                        cells.append(f"{column_name}{row_index}={value}")
                    formatted_rows.append(f"Row {row_index}: " + " | ".join(cells))
            return "\n".join(formatted_rows)
        except Exception as e:
            logger.error(f"Error parsing CSV file {file_path}: {e}")
            raise ValueError(f"Could not parse CSV content: {str(e)}")

    elif file_type == "xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=False, data_only=False)
            all_text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                if ws.max_row < 1 or ws.max_column < 1:
                    continue
                all_text_parts.append(f"Sheet: {sheet_name} | Used range: A1:{ws.cell(ws.max_row, ws.max_column).coordinate}")
                for row_index in range(1, ws.max_row + 1):
                    cells = []
                    for column_index in range(1, ws.max_column + 1):
                        cell = ws.cell(row_index, column_index)
                        value = "" if cell.value is None else str(cell.value)
                        cells.append(f"{cell.coordinate}={value}")
                    all_text_parts.append(f"Row {row_index}: " + " | ".join(cells))
            wb.close()
            return "\n".join(all_text_parts)
        except Exception as e:
            logger.error(f"Error parsing Excel file {file_path}: {e}")
            raise ValueError(f"Could not parse Excel content: {str(e)}")

    elif file_type == "docx":
        try:
            from docx import Document as WordDocument
            document = WordDocument(file_path)
            parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
            for table_index, table in enumerate(document.tables, start=1):
                parts.append(f"Table {table_index}:")
                for row_index, row in enumerate(table.rows, start=1):
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(f"Row {row_index}: " + " | ".join(cells))
            return "\n".join(parts)
        except Exception as e:
            logger.error(f"Error parsing Word file {file_path}: {e}")
            raise ValueError(f"Could not parse Word content: {str(e)}")

    elif file_type == "pptx":
        try:
            from pptx import Presentation
            presentation = Presentation(file_path)
            parts = []
            for slide_index, slide in enumerate(presentation.slides, start=1):
                slide_parts = [f"Slide {slide_index}:"]
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            slide_parts.append(" | ".join(cell.text.strip() for cell in row.cells))
                parts.append("\n".join(slide_parts))
            return "\n\n".join(parts)
        except Exception as e:
            logger.error(f"Error parsing PowerPoint file {file_path}: {e}")
            raise ValueError(f"Could not parse PowerPoint content: {str(e)}")
            
    elif file_type in ["mp3", "wav", "m4a", "ogg", "flac"]:
        try:
            try:
                transcript = _transcribe_local_media(file_path)
                return f"[Actual audio transcription for {os.path.basename(file_path)}]\n{transcript}"
            except ImportError:
                filename = os.path.basename(file_path)
                return f"[Audio File: {filename}] Audio content uploaded. Local faster-whisper package not installed. Size: {os.path.getsize(file_path)} bytes."
        except Exception as e:
            logger.error(f"Error parsing audio file {file_path}: {e}")
            raise ValueError(f"Could not parse audio content: {str(e)}")

    elif file_type in ["mp4", "avi", "mkv", "webm", "mov", "flv"]:
        try:
            try:
                transcript = _transcribe_local_media(file_path)
            except Exception as exc:
                logger.warning("Video speech transcription failed: %s", exc)
                transcript = ""
            try:
                from app.youtube_analysis import _frames, _caption_frames
                frames = _frames(file_path, count=8)
                visual = _caption_frames(frames) if frames else ""
            except Exception as exc:
                logger.warning("Video frame analysis failed: %s", exc)
                visual = ""
            evidence = [f"[Actual uploaded-video analysis for {os.path.basename(file_path)}]"]
            if transcript:
                evidence.append("Speech/transcript:\n" + transcript)
            if visual:
                evidence.append("Chronological sampled-frame descriptions:\n" + visual)
            if not transcript and not visual:
                evidence.append("No speech or visual evidence could be extracted; do not guess the contents.")
            return "\n\n".join(evidence)
        except Exception as e:
            logger.error(f"Error parsing video file {file_path}: {e}")
            raise ValueError(f"Could not parse video content: {str(e)}")

    elif file_type in ["png", "jpg", "jpeg", "webp", "bmp", "tiff"]:
        try:
            import base64
            import requests
            with open(file_path, "rb") as f:
                img_bytes = f.read()
            service_url = os.getenv("LOCAL_IMAGE_SERVICE_URL", "http://media-generator:8002")
            response = requests.post(
                f"{service_url}/caption",
                json={"frames": [base64.b64encode(img_bytes).decode("ascii")], "ocr": True},
                timeout=600,
            )
            response.raise_for_status()
            result = response.json()
            captions = result.get("captions", [])
            ocr_results = result.get("ocr", [])
            evidence = [f"[Actual local image analysis for {os.path.basename(file_path)}]"]
            if captions:
                evidence.append("Visual description:\n" + "\n".join(captions))
            if ocr_results:
                evidence.append("Extracted text (OCR):\n" + "\n".join(ocr_results))
            if not captions and not ocr_results:
                raise ValueError("The local image analyzer returned no visual or OCR evidence.")
            return "\n\n".join(evidence)
        except Exception as e:
            logger.error(f"Error calling local image analyzer: {e}")
            raise ValueError(
                f"[Image File: {os.path.basename(file_path)}] "
                "Image analysis is currently unavailable for this file type. "
                "The model cannot read this image directly. "
                "Please use a vision-capable model for image analysis."
            )

    elif file_type in ["txt", "md", "xml", "py", "cpp", "h", "json", "yaml", "yml", "log", "html", "htm"]:
        try:
            with open(file_path, mode="r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            if file_type == "json":
                # Beautify JSON for embedding readability
                try:
                    parsed = json.loads(content)
                    return json.dumps(parsed, indent=2)
                except Exception:
                    pass
            return content
        except Exception as e:
            logger.error(f"Error parsing text file {file_path}: {e}")
            raise ValueError(f"Could not read text content: {str(e)}")
            
    else:
        # Fallback raw read
        try:
            with open(file_path, mode="r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            raise ValueError(f"Unsupported file format: {file_type}. Parse failed: {str(e)}")


def fetch_url_content(url: str) -> str:
    """Fetch and extract clean readable text from a URL.
    
    Supports:
    - Regular web pages (via requests + BeautifulSoup)
    - YouTube URLs (extracts title + description)
    
    Returns extracted text for AI processing.
    """
    try:
        import re as _re
        # YouTube URL handling
        yt_patterns = [
            r'(?:youtube\.com/(?:watch\?[^\s]*?v=|shorts/|live/)|youtu\.be/)([\w-]{6,})',
        ]
        for pattern in yt_patterns:
            match = _re.search(pattern, url)
            if match:
                video_id = match.group(1)
                try:
                    from app.youtube_analysis import analyze_youtube_video
                    result = analyze_youtube_video(url, video_id)
                    return f"[YouTube Video]\nURL: {result['url']}\n\n{result['snippet']}"
                except Exception as exc:
                    return f"[YouTube Video: {video_id}]\nURL: {url}\nContent extraction failed; do not guess. Error: {exc}"

        # General web page extraction
        import requests as _r
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120.0.0.0 Safari/537.36'
        }
        renderer_url = os.getenv("BROWSER_RENDER_SERVICE_URL", "http://browser-renderer:8003")

        # Fast path for ordinary server-rendered pages.
        try:
            fast_resp = _safe_public_get(url, headers=headers, timeout=10)
            fast_resp.raise_for_status()
            fast_content_type = fast_resp.headers.get("content-type", "")
            if "text/html" in fast_content_type or not fast_content_type:
                from bs4 import BeautifulSoup
                fast_soup = BeautifulSoup(fast_resp.text, "html.parser")
                
                title = fast_soup.title.string.strip() if (fast_soup.title and fast_soup.title.string) else ""
                meta_desc = ""
                for m in fast_soup.find_all("meta"):
                    name = (m.get("name") or "").lower()
                    prop = (m.get("property") or "").lower()
                    content = (m.get("content") or "").strip()
                    if not title and prop == "og:title":
                        title = content
                    if name == "description" or prop in ("og:description", "twitter:description"):
                        if not meta_desc or len(content) > len(meta_desc):
                            meta_desc = content

                for tag in fast_soup(["script", "style", "nav", "footer", "header", "aside", "form", "iframe", "svg", "noscript"]):
                    tag.decompose()
                
                fast_main = fast_soup.find("article") or fast_soup.find("main") or fast_soup.find("body")
                fast_text = fast_main.get_text(separator="\n", strip=True) if fast_main else fast_soup.get_text(separator="\n", strip=True)
                lines = [line.strip() for line in fast_text.splitlines() if line.strip()]
                clean_body = "\n".join(lines)
                
                parts = [f"[Web Page: {title or urllib.parse.urlparse(url).netloc}]", f"URL: {fast_resp.url}"]
                if meta_desc:
                    parts.append(f"Meta Description: {meta_desc}")
                if clean_body:
                    parts.append(f"Content:\n{clean_body[:50000]}")
                
                final_excerpt = "\n\n".join(parts)
                if len(clean_body) >= 200 or meta_desc:
                    return final_excerpt
        except Exception:
            pass

        try:
            _validate_public_url(url)
            rendered = _r.post(f"{renderer_url}/render", json={"url": url}, timeout=40)
            if rendered.ok:
                rendered_data = rendered.json()
                rendered_text = rendered_data.get("text", "").strip()
                if rendered_text:
                    if len(rendered_text) > 50000:
                        rendered_text = rendered_text[:50000] + "\n\n[Content truncated...]"
                    return f"[Web Page: {rendered_data.get('title', '')}]\nURL: {rendered_data.get('url', url)}\n\n{rendered_text}"
        except Exception as browser_error:
            logger.warning("Local browser rendering failed, using HTML fallback: %s", browser_error)
        resp = _safe_public_get(url, headers=headers, timeout=15)
        resp.raise_for_status()
        content_type = resp.headers.get("content-type", "")
        
        if "text/html" in content_type or not content_type:
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(resp.text, "html.parser")
                title = soup.title.string.strip() if (soup.title and soup.title.string) else ""
                meta_desc = ""
                for m in soup.find_all("meta"):
                    name = (m.get("name") or "").lower()
                    prop = (m.get("property") or "").lower()
                    content = (m.get("content") or "").strip()
                    if not title and prop == "og:title":
                        title = content
                    if name == "description" or prop in ("og:description", "twitter:description"):
                        if not meta_desc or len(content) > len(meta_desc):
                            meta_desc = content
                for tag in soup(["script", "style", "nav", "footer", "header", "aside", "form", "iframe", "svg", "noscript"]):
                    tag.decompose()
                main = soup.find("article") or soup.find("main") or soup.find("body")
                text = main.get_text(separator="\n", strip=True) if main else soup.get_text(separator="\n", strip=True)
                lines = [l.strip() for l in text.splitlines() if l.strip()]
                clean_text = "\n".join(lines)
                
                parts = [f"[Web Page: {title or urllib.parse.urlparse(url).netloc}]", f"URL: {url}"]
                if meta_desc:
                    parts.append(f"Meta Description: {meta_desc}")
                if clean_text:
                    parts.append(f"Content:\n{clean_text[:50000]}")
                return "\n\n".join(parts)
            except ImportError:
                text = resp.text[:50000]
                return f"[Web Page]\nURL: {url}\n\n{text}"
        elif "application/json" in content_type:
            return f"[JSON API Response]\nURL: {url}\n\n{resp.text[:50000]}"
        elif "text/plain" in content_type:
            return f"[Text Content]\nURL: {url}\n\n{resp.text[:50000]}"
        else:
            return f"[Content from URL: {url}]\nContent-Type: {content_type}\n\n{resp.text[:10000]}"
    except Exception as e:
        raise ValueError(f"Could not fetch URL content from '{url}': {str(e)}")


# Helper processes (nvidia-smi, wmic, espeak, ...) are polled frequently. On
# Windows each one would otherwise flash a console window on screen, which in a
# windowed desktop build looks like terminals opening and closing by themselves.
_NO_WINDOW = getattr(subprocess, "CREATE_NO_WINDOW", 0)


def _run_hidden(command, **kwargs):
    """Run a helper process without ever showing a console window."""
    kwargs.setdefault("capture_output", True)
    kwargs.setdefault("text", True)
    if _NO_WINDOW:
        kwargs["creationflags"] = kwargs.get("creationflags", 0) | _NO_WINDOW
    return subprocess.run(command, **kwargs)


def _detect_cpu_name_windows() -> str:
    """Best-effort CPU name detection on Windows using WMIC or PowerShell."""
    name = ""
    try:
        res = _run_hidden(
            ["wmic", "cpu", "get", "name"],
            capture_output=True, text=True, timeout=5
        )
        if res.returncode == 0:
            lines = [ln.strip() for ln in res.stdout.strip().splitlines() if ln.strip()]
            if len(lines) >= 2:
                name = lines[1]
    except Exception:
        pass
    if not name:
        try:
            res = _run_hidden(
                ["powershell", "-Command", "Get-WmiObject Win32_Processor | Select-Object -ExpandProperty Name"],
                capture_output=True, text=True, timeout=5
            )
            if res.returncode == 0:
                name = res.stdout.strip()
        except Exception:
            pass
    return name or ""


def _detect_gpu_name_windows() -> str:
    """Best-effort GPU name detection on Windows using WMIC or PowerShell."""
    name = ""
    try:
        res = _run_hidden(
            ["wmic", "path", "win32_VideoController", "get", "name"],
            capture_output=True, text=True, timeout=5
        )
        if res.returncode == 0:
            lines = [ln.strip() for ln in res.stdout.strip().splitlines() if ln.strip()]
            if len(lines) >= 2:
                name = lines[1]
    except Exception:
        pass
    if not name:
        try:
            res = _run_hidden(
                ["powershell", "-Command", "Get-WmiObject Win32_VideoController | Select-Object -ExpandProperty Name"],
                capture_output=True, text=True, timeout=5
            )
            if res.returncode == 0:
                name = res.stdout.strip()
        except Exception:
            pass
    return name or ""


# --- System Performance Diagnostics ---
def get_folder_size_mb(path: str) -> float:
    """Return the total size of files under the specified path in MB."""
    if not os.path.exists(path):
        return 0.0
    if os.path.isfile(path):
        return os.path.getsize(path) / (1024 * 1024)
        
    total_size = 0
    for dirpath, _, filenames in os.walk(path):
        for f in filenames:
            fp = os.path.join(dirpath, f)
            if os.path.exists(fp):
                total_size += os.path.getsize(fp)
    return total_size / (1024 * 1024)

import time

# Performance trackers for rate calculations (disk read/write speed and network traffic speed)
_last_telemetry_time = None
_last_net_io = None
_last_disk_io = None
_cached_gpu_info = None
_last_gpu_check = 0

# ---------- REAL CPU NAME DETECTION (WMI on Windows, /proc on Linux) ----------
_cached_cpu_name = None
_cached_cpu_cores = None
_cached_cpu_threads = None
_cached_system_info = None

def _detect_real_cpu():
    """Detect the real CPU name using WMI (Windows) or /proc/cpuinfo (Linux).
    Caches result so WMI is only called once per process lifetime."""
    global _cached_cpu_name, _cached_cpu_cores, _cached_cpu_threads, _cached_system_info
    if _cached_cpu_name is not None:
        return _cached_cpu_name, _cached_cpu_cores, _cached_cpu_threads

    is_windows = platform.system() == "Windows"

    if is_windows:
        # WMI gives the exact retail CPU name, e.g. "AMD Ryzen 9 4900H with Radeon Graphics"
        try:
            r = _run_hidden(
                ["powershell", "-NoProfile", "-Command",
                 "Get-CimInstance -ClassName Win32_Processor | Select-Object Name, NumberOfCores, NumberOfLogicalProcessors | ConvertTo-Json"],
                capture_output=True, text=True, timeout=10
            )
            if r.returncode == 0 and r.stdout.strip():
                cpu = json.loads(r.stdout.strip())
                if isinstance(cpu, list):
                    cpu = cpu[0]
                _cached_cpu_name = (cpu.get("Name") or "").strip()
                _cached_cpu_cores = cpu.get("NumberOfCores", 0) or 0
                _cached_cpu_threads = cpu.get("NumberOfLogicalProcessors", 0) or 0
        except Exception:
            pass

        # Also grab system manufacturer/model
        try:
            r = _run_hidden(
                ["powershell", "-NoProfile", "-Command",
                 "Get-CimInstance -ClassName Win32_ComputerSystem | Select-Object Manufacturer, Model, TotalPhysicalMemory | ConvertTo-Json"],
                capture_output=True, text=True, timeout=10
            )
            if r.returncode == 0 and r.stdout.strip():
                _cached_system_info = json.loads(r.stdout.strip())
        except Exception:
            pass
    else:
        # Linux: /proc/cpuinfo
        try:
            with open("/proc/cpuinfo") as f:
                for line in f:
                    if "model name" in line:
                        _cached_cpu_name = line.split(":")[1].strip()
                        break
        except Exception:
            pass
        _cached_cpu_cores = psutil.cpu_count(logical=False) or 0
        _cached_cpu_threads = psutil.cpu_count(logical=True) or 0

    # Fallback: platform.processor() is better than "Unknown"
    if not _cached_cpu_name or _cached_cpu_name == "Unknown":
        proc = platform.processor()
        if proc:
            _cached_cpu_name = proc
        else:
            _cached_cpu_name = "Unknown CPU"

    if not _cached_cpu_cores:
        _cached_cpu_cores = psutil.cpu_count(logical=False) or 4
    if not _cached_cpu_threads:
        _cached_cpu_threads = psutil.cpu_count(logical=True) or 8

    logger.info(f"Detected CPU: {_cached_cpu_name} ({_cached_cpu_cores}C/{_cached_cpu_threads}T)")
    return _cached_cpu_name, _cached_cpu_cores, _cached_cpu_threads

# ---------- REAL GPU DETECTION (nvidia-smi + WMI fallback for integrated GPUs) ----------
_cached_wmi_gpus = None

def _detect_wmi_gpus():
    """Detect GPUs via WMI on Windows. Returns list of dicts with name and vram."""
    global _cached_wmi_gpus
    if _cached_wmi_gpus is not None:
        return _cached_wmi_gpus
    _cached_wmi_gpus = []
    if platform.system() != "Windows":
        return _cached_wmi_gpus
    try:
        r = _run_hidden(
            ["powershell", "-NoProfile", "-Command",
             "Get-CimInstance -ClassName Win32_VideoController | Select-Object Name, AdapterRAM, DriverVersion | ConvertTo-Json"],
            capture_output=True, text=True, timeout=10
        )
        if r.returncode == 0 and r.stdout.strip():
            gpus = json.loads(r.stdout.strip())
            if isinstance(gpus, dict):
                gpus = [gpus]
            for g in gpus:
                name = (g.get("Name") or "").strip()
                adapter_ram = g.get("AdapterRAM", 0) or 0
                vram_gb = round(adapter_ram / (1024**3), 2) if adapter_ram > 0 else 0
                _cached_wmi_gpus.append({
                    "name": name,
                    "vram_total_gb": vram_gb,
                    "driver": g.get("DriverVersion", ""),
                })
            logger.info(f"WMI detected {len(_cached_wmi_gpus)} GPU(s): {[g['name'] for g in _cached_wmi_gpus]}")
    except Exception as e:
        logger.warning(f"WMI GPU detection failed: {e}")
    return _cached_wmi_gpus

# ---------- REAL RAM DETECTION ----------
_cached_ram_total_gb = None

def _detect_real_ram():
    """Get accurate total RAM from WMI on Windows, psutil on Linux."""
    global _cached_ram_total_gb
    if _cached_ram_total_gb is not None:
        return _cached_ram_total_gb
    if _cached_system_info and _cached_system_info.get("TotalPhysicalMemory"):
        _cached_ram_total_gb = round(int(_cached_system_info["TotalPhysicalMemory"]) / (1024**3), 2)
    else:
        try:
            mem = psutil.virtual_memory()
            _cached_ram_total_gb = round(mem.total / (1024**3), 2)
        except Exception:
            _cached_ram_total_gb = 0
    return _cached_ram_total_gb

# ---------- TOKEN/SEC AND RESPONSE TIME TRACKING ----------
_inference_stats = {
    "total_tokens": 0,
    "total_requests": 0,
    "total_time_sec": 0.0,
    "last_tokens_per_sec": 0.0,
    "last_response_time_ms": 0.0,
    "avg_tokens_per_sec": 0.0,
    "avg_response_time_ms": 0.0,
    "token_measurement_source": "unavailable",
}

def record_inference_metrics(tokens_generated: int, elapsed_sec: float, source: str = "runtime_reported"):
    """Track only token counts reported by a live inference runtime."""
    global _inference_stats
    if elapsed_sec > 0 and tokens_generated > 0:
        tps = tokens_generated / elapsed_sec
        _inference_stats["last_tokens_per_sec"] = round(tps, 1)
        _inference_stats["last_response_time_ms"] = round(elapsed_sec * 1000, 0)
        _inference_stats["total_tokens"] += tokens_generated
        _inference_stats["total_requests"] += 1
        _inference_stats["total_time_sec"] += elapsed_sec
        _inference_stats["token_measurement_source"] = source or "runtime_reported"
        if _inference_stats["total_time_sec"] > 0:
            _inference_stats["avg_tokens_per_sec"] = round(
                _inference_stats["total_tokens"] / _inference_stats["total_time_sec"], 1
            )
        if _inference_stats["total_requests"] > 0:
            _inference_stats["avg_response_time_ms"] = round(
                (_inference_stats["total_time_sec"] * 1000) / _inference_stats["total_requests"], 0
            )

# How long a GPU reading may be reused. When nvidia-smi supplies live
# utilisation, a long cache freezes the number (it looked stuck at one value and
# drew a flat graph), so live sources are re-read almost every poll. Static
# identity-only results can be cached for much longer.
_GPU_LIVE_CACHE_SECONDS = 2
_GPU_STATIC_CACHE_SECONDS = 60


def _get_static_gpus():
    global _cached_gpu_info, _last_gpu_check
    now = time.time()
    if _cached_gpu_info is not None:
        has_live = any(gpu.get("has_live_metrics") for gpu in _cached_gpu_info)
        ttl = _GPU_LIVE_CACHE_SECONDS if has_live else _GPU_STATIC_CACHE_SECONDS
        if now - _last_gpu_check < ttl:
            return _cached_gpu_info

    gpus = []
    _last_gpu_check = now

    # Try nvidia-smi first (gives live metrics)
    try:
        res = _run_hidden(
            ["nvidia-smi", "--query-gpu=index,name,memory.total,memory.used,temperature.gpu,utilization.gpu", "--format=csv,noheader,nounits"],
            capture_output=True, text=True, timeout=2
        )
        if res.returncode == 0 and res.stdout.strip():
            for line in res.stdout.strip().split("\n"):
                parts = [p.strip() for p in line.split(",")]
                if len(parts) >= 6:
                    gpus.append({
                        "index": int(parts[0]),
                        "name": parts[1],
                        "vram_total_gb": round(float(parts[2]) / 1024.0, 2),
                        "vram_used_gb": round(float(parts[3]) / 1024.0, 2),
                        "temperature": float(parts[4]),
                        "usage": float(parts[5]),
                        "vendor": "nvidia",
                        "has_live_metrics": True
                    })
    except Exception:
        pass

    # If nvidia-smi found nothing, check WMI / lspci / Torch / dynamic detection
    if not gpus:
        # Check PyTorch CUDA availability
        try:
            import torch
            if torch.cuda.is_available():
                for idx in range(torch.cuda.device_count()):
                    props = torch.cuda.get_device_properties(idx)
                    vram_gb = round(props.total_memory / (1024**3), 2)
                    gpus.append({
                        "index": idx,
                        "name": props.name,
                        "vram_total_gb": vram_gb,
                        "vram_used_gb": round(torch.cuda.memory_allocated(idx) / (1024**3), 2),
                        "temperature": 0.0,
                        "usage": 0.0,
                        "vendor": "nvidia" if "nvidia" in props.name.lower() else "amd",
                        "has_live_metrics": True
                    })
        except Exception:
            pass

        # Check WMI GPUs (Windows native or host)
        if not gpus and platform.system() == "Windows":
            wmi_gpus = _detect_wmi_gpus()
            for idx, wg in enumerate(wmi_gpus):
                if wg["name"] and "Microsoft Basic" not in wg["name"]:
                    gpus.append({
                        "index": idx,
                        "name": wg["name"],
                        "vram_total_gb": wg["vram_total_gb"],
                        "vram_used_gb": 0.0,
                        "temperature": 0.0,
                        "usage": 0.0,
                        "vendor": "amd" if "amd" in wg["name"].lower() or "radeon" in wg["name"].lower() else "intel" if "intel" in wg["name"].lower() else "nvidia" if "nvidia" in wg["name"].lower() else "unknown",
                        "has_live_metrics": False
                    })

        # Check hardware_config.json if written by host bridge
        if not gpus:
            candidate_paths = [
                os.path.join(os.getenv("DATA_DIR", "/app/data"), "hardware_config.json"),
                os.path.join("/app", "hardware_config.json"),
                os.path.join(os.path.dirname(__file__), "..", "hardware_config.json"),
            ]
            for hw_file in candidate_paths:
                if os.path.exists(hw_file):
                    try:
                        with open(hw_file, "r") as f:
                            hw_data = json.load(f)
                            gpu_list = hw_data.get("gpus") or hw_data.get("host_gpus")
                            if gpu_list and isinstance(gpu_list, list) and len(gpu_list) > 0:
                                for idx, g in enumerate(gpu_list):
                                    g_name = g.get("name") or "Graphics Adapter"
                                    g_vram = float(g.get("vram_gb") or g.get("vram_total_gb") or 0.0)
                                    gpus.append({
                                        "index": idx,
                                        "name": g_name,
                                        "vram_total_gb": g_vram,
                                        "vram_used_gb": 0.0,
                                        "temperature": 0.0,
                                        "usage": 0.0,
                                        "vendor": "nvidia" if "nvidia" in g_name.lower() else ("amd" if "amd" in g_name.lower() or "radeon" in g_name.lower() else "intel"),
                                        "has_live_metrics": True
                                    })
                                break
                            elif hw_data.get("gpu_name") and hw_data["gpu_name"] not in ["N/A", "No GPU", ""]:
                                g_name = hw_data["gpu_name"]
                                g_vram = float(hw_data.get("gpu_vram_gb", 0.0) or 0.0)
                                gpus.append({
                                    "index": 0,
                                    "name": g_name,
                                    "vram_total_gb": g_vram,
                                    "vram_used_gb": 0.0,
                                    "temperature": 0.0,
                                    "usage": 0.0,
                                    "vendor": "nvidia" if "nvidia" in g_name.lower() else ("amd" if "amd" in g_name.lower() or "radeon" in g_name.lower() else "intel"),
                                    "has_live_metrics": True
                                })
                                break
                            elif hw_data.get("host_gpu_name") and hw_data["host_gpu_name"] not in ["N/A", "No GPU", ""]:
                                g_name = hw_data["host_gpu_name"]
                                g_vram = float(hw_data.get("host_gpu_vram_gb", 0.0) or 0.0)
                                gpus.append({
                                    "index": 0,
                                    "name": g_name,
                                    "vram_total_gb": g_vram,
                                    "vram_used_gb": 0.0,
                                    "temperature": 0.0,
                                    "usage": 0.0,
                                    "vendor": "nvidia" if "nvidia" in g_name.lower() else ("amd" if "amd" in g_name.lower() or "radeon" in g_name.lower() else "intel"),
                                    "has_live_metrics": True
                                })
                                break
                    except Exception:
                        pass
    # If we got nvidia GPUs but also have integrated, merge them
    elif platform.system() == "Windows":
        wmi_gpus = _detect_wmi_gpus()
        nvidia_names = {g["name"].lower() for g in gpus}
        for idx, wg in enumerate(wmi_gpus):
            if wg["name"] and "Microsoft" not in wg["name"] and wg["name"].lower() not in nvidia_names:
                # Check if it's not a duplicate of an existing nvidia GPU
                is_dupe = any(wg["name"].lower() in n for n in nvidia_names)
                if not is_dupe:
                    gpus.append({
                        "index": len(gpus),
                        "name": wg["name"],
                        "vram_total_gb": wg["vram_total_gb"],
                        "vram_used_gb": 0,
                        "temperature": 0,
                        "usage": 0,
                        "vendor": "amd" if "amd" in wg["name"].lower() or "radeon" in wg["name"].lower() else "intel",
                        "has_live_metrics": False
                    })

    _cached_gpu_info = gpus
    return gpus

def get_system_telemetry(db: Session = None, active_sessions: int = 0, latency_ms: float = 0.0, **_ignored) -> dict:
    """Calculate system-wide resource metrics using auto-detected host hardware specs.
    
    Priority:
       1. host_stats.json - written every second by the optional Windows/macOS/Linux
                            host bridge. Values are source-labelled measurements.
       2. hardware_config.json - written by bootstrapper.py at container startup.
       3. /host/proc/* - host procfs mounted into container (Docker --privileged or volume mount).
       4. Environment variables - SMARAN_HOST_CPU_NAME, SMARAN_HOST_RAM_GB, SMARAN_HOST_GPU_NAME, etc.
       5. nvidia-smi / torch.cuda - GPU passthrough inside Docker.
       6. psutil - container-level fallback (RAM/CPU will reflect container limits).
    """
    global _last_telemetry_time, _last_net_io, _last_disk_io

    now = time.time()
    dt = now - _last_telemetry_time if _last_telemetry_time else 1.0
    if dt <= 0:
        dt = 1.0
    _last_telemetry_time = now

    # Try reading from host_stats_bridge output (most accurate)
    _hs = {}
    try:
        data_dir = os.getenv("DATA_DIR", "/app/data")
        hs_path = os.path.join(data_dir, "host_stats.json")
        if os.path.exists(hs_path):
            age = time.time() - os.path.getmtime(hs_path)
            if age < 5:
                with open(hs_path) as f:
                    _hs = json.load(f)
    except Exception:
        pass

    # Also read hardware_config.json for static hardware specs
    _hw = {}
    try:
        hw_path = os.path.join(os.getenv("DATA_DIR", "/app/data"), "hardware_config.json")
        if os.path.exists(hw_path):
            with open(hw_path) as f:
                _hw = json.load(f)
    except Exception:
        pass

    # Environment variable overrides (Docker -e flags or docker-compose env)
    env_cpu_name = os.getenv("SMARAN_HOST_CPU_NAME", "").strip()
    env_ram_gb = os.getenv("SMARAN_HOST_RAM_GB", "").strip()

    # 1. CPU
    if _hs:
        cpu_usage  = float(_hs.get("cpu_usage", _hs.get("cpu_usage_percent", 0.0)))
        cpu_name   = str(_hs.get("cpu_name", ""))
        cpu_cores  = int(_hs.get("cpu_cores", _hs.get("cpu_physical_cores", 0)))
        cpu_threads = int(_hs.get("cpu_threads", _hs.get("cpu_logical_threads", 0)))
    else:
        cpu_usage = psutil.cpu_percent(interval=None)
        real_cpu_name, real_cores, real_threads = _detect_real_cpu()
        cpu_name  = env_cpu_name or str(_hw.get("host_cpu_name", "")) or real_cpu_name
        if not cpu_name or cpu_name == "Unknown CPU":
            cpu_name = real_cpu_name
        cpu_cores = int(_hw.get("host_cpu_cores", 0)) or real_cores or psutil.cpu_count(logical=False) or 4
        cpu_threads = int(_hw.get("host_cpu_threads", 0)) or real_threads or psutil.cpu_count(logical=True) or 8

    # 2. Memory
    if _hs:
        mem_pct      = float(_hs.get("ram_percent", 0.0))
        mem_used_gb  = float(_hs.get("ram_used_gb", 0.0))
        mem_total_gb = float(_hs.get("ram_total_gb", 0.0))
    else:
        host_ram_total = float(env_ram_gb or _hw.get("host_ram_total_gb", 0) or _detect_real_ram() or 0)
        try:
            mem = psutil.virtual_memory()
            mem_pct = mem.percent
            mem_total_gb = host_ram_total if host_ram_total > 0 else round(mem.total / (1024**3), 2)
            mem_used_gb  = round((mem_pct / 100.0) * mem_total_gb, 2)
        except Exception:
            mem_pct = mem_used_gb = mem_total_gb = 0.0

    # 3. GPU
    gpus = _hs.get("gpus", []) if _hs and isinstance(_hs.get("gpus"), list) else _get_static_gpus()
    gpu_available = bool(gpus)
    primary_gpu = next((gpu for gpu in gpus if gpu.get("has_live_metrics")), gpus[0] if gpus else {})
    gpu_name = primary_gpu.get("name", "N/A")
    gpu_vram_total = primary_gpu.get("vram_total_gb")
    gpu_vram_used = primary_gpu.get("vram_used_gb")
    gpu_temperature = primary_gpu.get("temperature")
    gpu_usage = primary_gpu.get("usage")

    # 4. Disk
    if _hs:
        disk_io_pct       = float(_hs.get("disk_io_pct", 0.0))
        disk_space_pct    = float(_hs.get("disk_space_pct", 0.0))
        disk_used_gb      = float(_hs.get("disk_space_used_gb", _hs.get("disk_used_gb", 0.0)))
        disk_total_gb     = float(_hs.get("disk_space_total_gb", _hs.get("disk_total_gb", 0.0)))
        disk_read_kb      = float(_hs.get("disk_read_kb", 0.0))
        disk_write_kb     = float(_hs.get("disk_write_kb", 0.0))
    else:
        disk_io_pct = disk_space_pct = disk_used_gb = disk_total_gb = 0.0
        disk_read_kb = disk_write_kb = 0.0
        try:
            disk = psutil.disk_usage('/')
            disk_space_pct = round(disk.percent, 1)
            disk_used_gb   = round(disk.used  / (1024**3), 2)
            disk_total_gb  = round(disk.total / (1024**3), 2)
        except Exception:
            pass

    # 5. Network
    if _hs:
        net_up_kb   = float(_hs.get("net_up_kb", 0.0))
        net_down_kb = float(_hs.get("net_down_kb", 0.0))
    else:
        net_up_kb = net_down_kb = 0.0
        try:
            net_io = psutil.net_io_counters()
            if net_io and _last_net_io and dt > 0:
                net_up_kb   = round(((net_io.bytes_sent - _last_net_io.bytes_sent) / 1024.0) / dt, 1)
                net_down_kb = round(((net_io.bytes_recv - _last_net_io.bytes_recv) / 1024.0) / dt, 1)
            _last_net_io = net_io
        except Exception:
            pass

    # 6. Database size
    sqlite_size = get_folder_size_mb(settings.SQLITE_DB_PATH)
    chroma_size = get_folder_size_mb(settings.CHROMA_DIR)
    db_size     = round(sqlite_size + chroma_size, 2)

    # 7. Model info
    model_display_name = _hw.get("display_name", "")
    ctx_window         = int(_hw.get("ctx_window", 0) or 0)
    reasoning_model    = bool(_hw.get("reasoning_model", False))
    selected_model_id  = str(_hw.get("model_id", ""))

    # Battery. Read fresh on every poll rather than cached: the browser's
    # Battery Status API is missing in the packaged desktop window and, where
    # it does exist, only reports at page load unless events are wired up —
    # which is why the charge level and the plugged-in symbol went stale.
    battery_percent = None
    battery_charging = None
    battery_minutes_left = None
    battery_present = False
    try:
        battery = psutil.sensors_battery()
    except Exception:
        battery = None
    if battery is not None:
        battery_present = True
        battery_percent = round(float(battery.percent), 1)
        battery_charging = bool(battery.power_plugged)
        seconds_left = getattr(battery, "secsleft", None)
        # psutil reports POWER_TIME_UNLIMITED/UNKNOWN as negative sentinels.
        if isinstance(seconds_left, int) and seconds_left >= 0:
            battery_minutes_left = seconds_left // 60

    return {
        "cpu_usage":          cpu_usage,
        "cpu_cores":          cpu_cores,
        "cpu_threads":        cpu_threads,
        "cpu_name":           cpu_name,
        "memory_usage":       mem_pct,
        "memory_used_gb":     mem_used_gb,
        "memory_total_gb":    mem_total_gb,
        "gpu_available":      gpu_available,
        "gpu_usage":          gpu_usage,
        "gpu_name":           gpu_name,
        "gpu_vram_used":      gpu_vram_used,
        "gpu_vram_total":     gpu_vram_total,
        "gpu_temperature":    gpu_temperature,
        "gpus":               gpus,
        "gpu_count":          len(gpus),
        "disk_usage":         disk_io_pct,
        "disk_space_pct":     disk_space_pct,
        "disk_used_gb":       disk_used_gb,
        "disk_total_gb":      disk_total_gb,
        "disk_read_kb":       max(0.0, disk_read_kb),
        "disk_write_kb":      max(0.0, disk_write_kb),
        "net_up_kb":          net_up_kb,
        "net_down_kb":        net_down_kb,
        "active_sessions":    active_sessions,
        "database_size_mb":   db_size,
        "average_latency_ms": round(latency_ms, 1),
        "tokens_per_sec":     _inference_stats.get("last_tokens_per_sec", 0.0),
        "avg_tokens_per_sec": _inference_stats.get("avg_tokens_per_sec", 0.0),
        "response_time_ms":   _inference_stats.get("last_response_time_ms", round(latency_ms, 1)),
        "avg_response_time_ms": _inference_stats.get("avg_response_time_ms", 0.0),
        "total_tokens":       _inference_stats.get("total_tokens", 0),
        "token_measurement_source": _inference_stats.get("token_measurement_source", "unavailable"),
        "model_display_name": model_display_name,
        "model_id":           selected_model_id,
        "ctx_window":         ctx_window,
        "reasoning_model":    reasoning_model,
        "telemetry_source":   str(_hs.get("telemetry_source") or "host_bridge") if _hs else "native_runtime",
        "host_os":            str(_hs.get("host_os") or "") if _hs else "",
        "host_os_display":    str(_hs.get("host_os_display") or "") if _hs else "",
        "host_os_version":    str(_hs.get("host_os_version") or "") if _hs else "",
        "host_os_build":      str(_hs.get("host_os_build") or "") if _hs else "",
        "host_arch":          str(_hs.get("host_arch") or "") if _hs else "",
        "host_device_manufacturer": str(_hs.get("host_device_manufacturer") or "") if _hs else "",
        "host_device_model":  str(_hs.get("host_device_model") or "") if _hs else "",
        "npu_available":      bool(_hs.get("npu_available", False)) if _hs else False,
        "npu_name":           str(_hs.get("npu_name") or "") if _hs else "",
        "npu_vendor":         str(_hs.get("npu_vendor") or "") if _hs else "",
        "battery_present":    battery_present,
        "battery_percent":    battery_percent,
        "battery_charging":   battery_charging,
        "battery_minutes_left": battery_minutes_left,
    }

import httpx

async def zep_add_message(session_id: str, role: str, content: str):
    """Asynchronously send chat messages to Zep AI Memory service."""
    zep_url = os.getenv("ZEP_URL", "http://zep-ai:8000")
    zep_role = "ai" if role == "assistant" else role
    payload = {
        "messages": [
            {
                "role": zep_role,
                "content": content
            }
        ]
    }
    try:
        async with httpx.AsyncClient(timeout=2.0) as client:
            url = f"{zep_url.rstrip('/')}/api/v1/sessions/{session_id}/memory"
            r = await client.post(url, json=payload)
            if r.status_code != 200:
                logger.warning(f"Zep AI memory add failed with status {r.status_code}: {r.text}")
    except Exception as e:
        logger.warning(f"Failed to connect to Zep AI for session {session_id}: {e}")

async def zep_get_history(session_id: str) -> list[dict]:
    """Retrieve sliding window memory history from Zep AI."""
    zep_url = os.getenv("ZEP_URL", "http://zep-ai:8000")
    try:
        async with httpx.AsyncClient(timeout=2.0) as client:
            url = f"{zep_url.rstrip('/')}/api/v1/sessions/{session_id}/memory"
            r = await client.get(url)
            if r.status_code == 200:
                data = r.json()
                messages = data.get("messages", [])
                history = []
                for msg in messages:
                    role = "assistant" if msg.get("role") == "ai" else msg.get("role")
                    history.append({"role": role, "content": msg.get("content", "")})
                return history
    except Exception as e:
        logger.warning(f"Failed to fetch history from Zep AI for session {session_id}: {e}")
    return []
